"""SCGCPR — Servicio de generación de exámenes con IA (Claude)."""
from __future__ import annotations
import json
from loguru import logger
from app.core.config import settings

# ---------------------------------------------------------------------------
# NOTE ON n_multi / n_casos PERSISTENCE
# ---------------------------------------------------------------------------
# FuenteIA has no dedicated columns for n_multi / n_casos / texto_pegado.
# Rather than add a migration, we serialize those values as JSON into the
# `prompt_usado` field when the endpoint creates the FuenteIA row:
#   prompt_usado = json.dumps({"n_multi": n, "n_casos": m, "texto_pegado": t})
# procesar_generacion_ia reads and parses that JSON to recover the values.
# If `ruta_archivo` is present, text is extracted from disk; otherwise
# `texto_pegado` from the JSON is used directly.
# ---------------------------------------------------------------------------


def extraer_texto_fuente(ruta: str, tipo_archivo: str) -> str:
    """Extrae texto de un archivo según su tipo (pdf|docx|pptx|texto)."""
    tipo = (tipo_archivo or "").lower()
    if tipo == "pdf":
        import pdfplumber
        with pdfplumber.open(ruta) as pdf:
            return "\n".join((page.extract_text() or "") for page in pdf.pages)
    if tipo == "docx":
        import docx
        d = docx.Document(ruta)
        return "\n".join(p.text for p in d.paragraphs)
    if tipo == "pptx":
        import pptx
        prs = pptx.Presentation(ruta)
        partes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    partes.append(shape.text_frame.text)
        return "\n".join(partes)
    if tipo == "texto":
        with open(ruta, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    raise ValueError(f"Tipo de archivo no soportado: {tipo_archivo}")


def construir_prompt(texto: str, n_multi: int, n_casos: int, producto: str | None = None,
                     n_vf: int = 0, n_objeciones: int = 0) -> str:
    """Construye el prompt base para la generación de preguntas con Claude.

    El texto de entrada se trunca a _TEXTO_MAX_CHARS para acotar el tamaño del
    prompt y el costo de tokens, independientemente del tamaño del archivo fuente.

    Las preguntas se formulan SOBRE el producto / la información (tratada como
    verdadera y válida), nunca "sobre el documento": la fuente se da por buena
    como ficha técnica o estudio clínico, y las preguntas deben sostenerse solas.
    """
    texto_acotado = texto[:_TEXTO_MAX_CHARS]
    total = n_multi + n_casos + n_vf + n_objeciones
    tema = (f" El material se relaciona con: {producto} (úsalo solo como tema/contexto; "
            "NO exijas que el material lo nombre explícitamente)." if producto else "")
    return (
        "Eres un experto en capacitación farmacéutica y en elaboración de exámenes. "
        "A partir del siguiente material de referencia —que debes tomar como VERDADERO "
        "y VÁLIDO— genera exactamente "
        f"{total} preguntas de evaluación profesional SOBRE EL CONTENIDO del material." + tema + "\n"
        f"- {n_multi} de opción múltiple (tipo 'multi', EXACTAMENTE 5 opciones a–e, 1 correcta)\n"
        f"- {n_casos} casos clínicos (tipo 'caso', con 'escenario', EXACTAMENTE 5 opciones a–e, 1 correcta)\n"
        f"- {n_vf} de Verdadero/Falso (tipo 'vf', EXACTAMENTE 2 opciones [\"Verdadero\",\"Falso\"] en ese orden; "
        "'correcta'=0 si la afirmación es verdadera, 1 si es falsa)\n"
        f"- {n_objeciones} de Objeción de Producto (tipo 'objecion'): en 'escenario' pon la OBJECIÓN "
        "que un médico plantea al visitador (una duda, efecto adverso, limitación o comparación con "
        "la competencia sobre el producto/tema); en 'texto' la pregunta '¿cuál es la mejor respuesta "
        "técnica?'; EXACTAMENTE 5 opciones a–e con 1 correcta (la mejor respuesta profesional)\n\n"
        "REGLAS DE REDACCIÓN (obligatorias):\n"
        "- Genera SIEMPRE las preguntas a partir del CONTENIDO real del material "
        "(conceptos, principios, indicaciones, dosis, contraindicaciones, efectos, datos, etc.). "
        "Si el material trata un tema general (p. ej. nutrición, fisiología, farmacología) y "
        "no menciona un producto de marca, formula las preguntas sobre ese TEMA/CONTENIDO.\n"
        "- NUNCA te niegues ni respondas con explicaciones, disculpas o comentarios fuera del "
        "JSON. Aunque el material no mencione el producto indicado, DEBES generar el examen a "
        "partir del contenido disponible y devolver el arreglo JSON de preguntas.\n"
        "- NUNCA menciones 'el documento', 'este documento', 'el texto', 'el material' "
        "ni la fuente. La pregunta debe entenderse por sí sola, como en un examen real.\n"
        "- Trata la información provista como un hecho establecido; no la cuestiones ni "
        "la relativices ('según el documento', 'el texto indica', etc. están prohibidos).\n\n"
        "Devuelve SOLO el arreglo JSON, sin envolverlo en bloques de código markdown "
        "(nada de ```), sin texto antes ni después. Cada pregunta con este esquema:\n"
        "tipo: 'multi'|'caso'|'vf'|'objecion'; escenario: string (obligatorio en caso y objecion); texto: string; "
        "opciones: array de strings (5 para multi/caso, 2 para vf); correcta: índice 0-based de la opción correcta; "
        "explicacion: string.\n\nMATERIAL DE REFERENCIA:\n" + texto_acotado)


# Maximum characters fed into the prompt to bound token cost.
_TEXTO_MAX_CHARS = 40_000


def _extraer_json(texto_respuesta: str):
    """Extrae y parsea JSON de la respuesta del modelo.

    Estrategia en orden de precedencia:
    1. Si hay un bloque ```json ... ``` o ``` ... ```, usa su contenido.
    2. De lo contrario busca el primer '[' o '{' y el último ']' o '}'
       correspondiente y parsea ese fragmento (tolerante a texto extra
       antes/después del JSON).
    3. Intenta parsear el texto completo como JSON.
    4. Lanza ValueError si todo falla.
    """
    t = texto_respuesta.strip()

    # Pre-limpieza: quitar fences markdown envolventes (```json … ``` o ``` … ```),
    # incluso si falta el cierre (respuesta truncada). Así las estrategias siguientes
    # ven el arreglo JSON directamente.
    if t.startswith("```"):
        import re as _re
        t = _re.sub(r"^```(?:json)?\s*", "", t)
        t = _re.sub(r"\s*```$", "", t).strip()

    # Strategy 1: fenced code block
    if "```" in t:
        import re
        # Match ```json ... ``` or ``` ... ```
        m = re.search(r"```(?:json)?\s*([\s\S]*?)```", t)
        if m:
            candidate = m.group(1).strip()
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                pass  # fall through to other strategies

    # Strategy 2: first balanced JSON array or object
    for open_ch, close_ch in (("[", "]"), ("{", "}")):
        start = t.find(open_ch)
        end = t.rfind(close_ch)
        if start != -1 and end != -1 and end > start:
            candidate = t[start:end + 1]
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                pass

    # Strategy 3: full text
    try:
        return json.loads(t)
    except json.JSONDecodeError as e:
        raise ValueError(f"La IA no devolvió JSON válido: {e}")


def _generar_demo(texto: str, n_multi: int, n_casos: int, producto: str | None = None,
                  n_vf: int = 0, n_objeciones: int = 0) -> list[dict]:
    """Genera preguntas localmente (MODO DEMO), sin llamar a Claude ni gastar API.

    Las preguntas se formulan SOBRE el producto / la información (tratada como
    verdadera y válida), nunca "sobre el documento". Se marca [DEMO] solo como
    indicador de modo. Pasan por validar_preguntas_generadas para validar el esquema.
    """
    import re
    frases = [s.strip() for s in re.split(r"[.\n]", texto) if len(s.strip()) >= 25]
    if not frases:
        frases = ["la información clínica del producto"]
    prod = (producto or "el producto").strip()
    # Distractores genéricos, plausibles y SIN referencia a la fuente/documento.
    distractores = [
        "No existe evidencia clínica que respalde esa afirmación.",
        "Está contraindicado en la totalidad de los pacientes.",
        "Carece de eficacia terapéutica demostrada.",
        "Corresponde a una indicación distinta a la establecida.",
    ]
    preguntas: list[dict] = []
    for i in range(max(0, n_multi)):
        frase = frases[i % len(frases)][:160]
        preguntas.append({
            "tipo": "multi",
            "texto": f"[DEMO] Respecto a {prod}, ¿cuál de las siguientes afirmaciones es correcta?",
            "opciones": [
                frase,
                distractores[i % len(distractores)],
                distractores[(i + 1) % len(distractores)],
                distractores[(i + 2) % len(distractores)],
                distractores[(i + 3) % len(distractores)],
            ],
            "correcta": 0,
            "explicacion": f"La afirmación correcta corresponde a la información clínica de {prod}. (Generada en MODO DEMO, sin IA real.)",
        })
    for i in range(max(0, n_vf)):
        frase = frases[(n_multi + i) % len(frases)][:170]
        preguntas.append({
            "tipo": "vf",
            "texto": f"[DEMO] {frase}.",   # afirmación verdadera tomada del contenido
            "opciones": ["Verdadero", "Falso"],
            "correcta": 0,
            "explicacion": f"Afirmación correcta según la información de {prod}. (Generada en MODO DEMO, sin IA real.)",
        })
    for i in range(max(0, n_casos)):
        frase = frases[(n_multi + n_vf + i) % len(frases)][:160]
        preguntas.append({
            "tipo": "caso",
            "escenario": (f"[DEMO] Un médico solicita información sobre {prod}. "
                          f"De acuerdo con la evidencia clínica disponible: {frase}."),
            "texto": f"¿Cuál es la información correcta que debe transmitir el representante sobre {prod}?",
            "opciones": [
                frase,
                distractores[i % len(distractores)],
                distractores[(i + 1) % len(distractores)],
                distractores[(i + 2) % len(distractores)],
                distractores[(i + 3) % len(distractores)],
            ],
            "correcta": 0,
            "explicacion": f"La respuesta correcta refleja la evidencia clínica de {prod}. (Generada en MODO DEMO, sin IA real.)",
        })
    for i in range(max(0, n_objeciones)):
        frase = frases[(n_multi + n_vf + n_casos + i) % len(frases)][:150]
        preguntas.append({
            "tipo": "objecion",
            "escenario": (f"[DEMO] Un médico objeta sobre {prod}: \"He escuchado dudas al respecto y "
                          f"prefiero no cambiar mi conducta\"."),
            "texto": f"¿Cuál es la mejor respuesta técnica del representante ante esta objeción sobre {prod}?",
            "opciones": [
                frase,
                distractores[i % len(distractores)],
                distractores[(i + 1) % len(distractores)],
                distractores[(i + 2) % len(distractores)],
                distractores[(i + 3) % len(distractores)],
            ],
            "correcta": 0,
            "explicacion": f"La mejor respuesta rebate la objeción con la evidencia clínica de {prod}. (Generada en MODO DEMO, sin IA real.)",
        })
    return validar_preguntas_generadas(preguntas)


def generar_preguntas_ia(texto: str, n_multi: int, n_casos: int,
                         client=None, model: str | None = None,
                         producto: str | None = None, n_vf: int = 0,
                         demo: bool | None = None, n_objeciones: int = 0) -> list[dict]:
    """Genera preguntas llamando a Claude. El cliente se inyecta para testing.

    `producto` ancla las preguntas al producto/tema (en vez de a "el documento").

    Si EXAMEN_IA_DEMO está activo y no se inyectó cliente, usa el generador local
    (sin consumir API) — permite probar el flujo completo sin créditos.
    """
    usar_demo = settings.EXAMEN_IA_DEMO if demo is None else demo
    if client is None and usar_demo:
        logger.info("generar_preguntas_ia: MODO DEMO activo — generación local sin Claude")
        return _generar_demo(texto, n_multi, n_casos, producto, n_vf, n_objeciones)
    if client is None:
        import anthropic
        client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
    model = model or settings.EXAM_AI_MODEL
    prompt = construir_prompt(texto, n_multi, n_casos, producto, n_vf, n_objeciones)
    # max_tokens proporcional a la cantidad pedida: cada pregunta en JSON ocupa
    # ~600-800 tokens. Con un tope fijo bajo (4096) los exámenes grandes se truncaban
    # y devolvían JSON inválido. Escalamos con un techo de seguridad.
    total = max(1, n_multi + n_casos + n_vf + n_objeciones)
    max_tokens = min(16000, 2000 + total * 800)
    respuesta = client.messages.create(
        model=model, max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
        # Desactiva el "extended thinking" (activo por defecto en modelos Claude
        # nuevos como sonnet-5): de lo contrario el razonamiento consume el
        # presupuesto de tokens y el bloque de texto (el JSON) puede llegar
        # vacío/truncado ("Expecting value: line 1 column 1"). Se pasa por
        # extra_body porque el SDK instalado no acepta el kwarg `thinking` nativo.
        extra_body={"thinking": {"type": "disabled"}})
    # Une solo el texto de bloques de tipo texto. Algunos modelos/SDK devuelven
    # bloques cuyo `.text` es None (p. ej. bloques de razonamiento): `or ""` evita
    # el TypeError de join ("expected str instance, NoneType found") y los ignora.
    texto_resp = "".join((getattr(b, "text", "") or "") for b in respuesta.content)
    # Red de seguridad: si el modelo responde en prosa (p. ej. negándose porque el
    # material no coincide con el producto) en vez de JSON, damos un error claro y
    # accionable en lugar del críptico "Expecting value: line 1 column 1 (char 0)".
    if "[" not in texto_resp and "{" not in texto_resp:
        motivo = texto_resp.strip()[:300] or "(respuesta vacía)"
        raise ValueError(
            "La IA no generó preguntas a partir del material. Respondió: "
            f"«{motivo}». Revisa que el documento tenga contenido suficiente, o "
            "ajusta/vacía el campo Producto para que coincida con el tema del material.")
    data = _extraer_json(texto_resp)
    return validar_preguntas_generadas(data)


def validar_preguntas_generadas(data: list) -> list[dict]:
    if not isinstance(data, list) or not data:
        raise ValueError("La IA no devolvió una lista de preguntas")
    out = []
    for i, q in enumerate(data):
        if not isinstance(q, dict):
            raise ValueError(f"Pregunta {i}: formato inválido")
        tipo = q.get("tipo")
        if tipo not in ("multi", "caso", "vf", "objecion"):
            raise ValueError(f"Pregunta {i}: tipo inválido {tipo!r}")
        if not (q.get("texto") or "").strip():
            raise ValueError(f"Pregunta {i}: texto vacío")
        ops = q.get("opciones")
        n_ops = 2 if tipo == "vf" else 5  # vf: Verdadero/Falso; multi/caso: 5 opciones (a–e)
        if not isinstance(ops, list) or len(ops) != n_ops or not all(isinstance(o, str) and o.strip() for o in ops):
            raise ValueError(f"Pregunta {i}: debe tener exactamente {n_ops} opciones no vacías")
        correcta = q.get("correcta")
        if not isinstance(correcta, int) or isinstance(correcta, bool) or not (0 <= correcta < n_ops):
            raise ValueError(f"Pregunta {i}: 'correcta' debe ser 0..{n_ops - 1}")
        if tipo == "caso" and not (q.get("escenario") or "").strip():
            raise ValueError(f"Pregunta {i}: caso clínico requiere 'escenario'")
        if tipo == "objecion" and not (q.get("escenario") or "").strip():
            raise ValueError(f"Pregunta {i}: la objeción requiere 'escenario' (el texto de la objeción)")
        out.append({
            "tipo": tipo, "texto": q["texto"].strip(),
            "escenario": (q.get("escenario") or None),
            "opciones": [o.strip() for o in ops], "correcta": correcta,
            "explicacion": (q.get("explicacion") or "").strip()})
    return out


# ---------------------------------------------------------------------------
# Persistence
# ---------------------------------------------------------------------------

def persistir_preguntas(db, examen_id: int, preguntas: list[dict]) -> int:
    """Inserta preguntas (con 4 opciones cada una) en el examen borrador.

    `indice_original` es 0..3 tal como las devolvió la IA.
    `es_correcta` se marca solo en la opción cuyo índice coincide con `correcta`.
    Retorna la cantidad de preguntas insertadas.
    """
    from app.models.exam_models import Pregunta, PreguntaOpcion

    base = db.query(Pregunta).filter(Pregunta.examen_id == examen_id).count()
    for offset, q in enumerate(preguntas):
        pregunta = Pregunta(
            examen_id=examen_id,
            tipo=q["tipo"],
            escenario=q.get("escenario"),
            texto=q["texto"],
            explicacion=q.get("explicacion"),
            orden=base + offset,
        )
        for idx, texto_op in enumerate(q["opciones"]):
            pregunta.opciones.append(
                PreguntaOpcion(
                    texto_opcion=texto_op,
                    indice_original=idx,
                    es_correcta=(idx == q["correcta"]),
                )
            )
        db.add(pregunta)
    db.commit()
    return len(preguntas)


# ---------------------------------------------------------------------------
# Background job
# ---------------------------------------------------------------------------

def procesar_generacion_ia(fuente_id: int) -> None:
    """Job en background: extrae texto, genera preguntas con IA y las persiste.

    Los parámetros n_multi / n_casos / texto_pegado se leen del JSON almacenado
    en FuenteIA.prompt_usado (no existe columna propia para ellos — decisión de
    diseño documentada en el módulo para evitar una migración adicional).
    """
    from app.db.database import SessionLocal
    from app.models.exam_models import Examen, FuenteIA

    db = SessionLocal()
    try:
        fuente = db.query(FuenteIA).filter(FuenteIA.id == fuente_id).first()
        if fuente is None:
            logger.warning(f"procesar_generacion_ia: FuenteIA id={fuente_id} no encontrada")
            return

        fuente.estado_generacion = "procesando"
        db.commit()

        try:
            # Recover parameters stored as JSON in prompt_usado
            params: dict = {}
            if fuente.prompt_usado:
                try:
                    params = json.loads(fuente.prompt_usado)
                except (json.JSONDecodeError, TypeError):
                    params = {}

            n_multi: int = int(params.get("n_multi", 5))
            n_casos: int = int(params.get("n_casos", 0))
            n_vf: int = int(params.get("n_vf", 0))
            n_objeciones: int = int(params.get("n_objeciones", 0))
            texto_pegado: str | None = params.get("texto_pegado") or None

            # Extract text: from file on disk, or from the pasted text
            if fuente.ruta_archivo:
                texto = extraer_texto_fuente(fuente.ruta_archivo, fuente.tipo_archivo or "texto")
            elif texto_pegado:
                texto = texto_pegado
            else:
                raise ValueError("No hay fuente de texto: ni archivo ni texto_pegado")

            examen = db.query(Examen).filter(Examen.id == fuente.examen_id).first()
            producto = examen.producto if examen else None
            # Flag DEMO en runtime (Config.DIM_Parametro), con fallback al .env.
            from app.services import config_service
            demo = config_service.obtener_bool(db, "EXAMEN_IA_DEMO", settings.EXAMEN_IA_DEMO)
            preguntas = generar_preguntas_ia(texto, n_multi, n_casos, producto=producto, n_vf=n_vf,
                                              demo=demo, n_objeciones=n_objeciones)
            persistir_preguntas(db, fuente.examen_id, preguntas)

            fuente.estado_generacion = "exitoso"
            fuente.mensaje_error = None

        except Exception as exc:
            fuente.estado_generacion = "error"
            fuente.mensaje_error = str(exc)[:1000]
            logger.error(f"procesar_generacion_ia fuente={fuente_id} falló: {exc}")

        db.commit()
    finally:
        db.close()
